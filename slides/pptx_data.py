from typing import Union
import textwrap
from dataclasses import dataclass, field
from pptx.dml.color import ColorFormat
from pptx.dml.fill import FillFormat
from pptx.util import Pt
from pptx.text.text import Font
from pptx.util import Length
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pygments import highlight

from code_formatting import PowerPointCodeFormatter, format_code_as_image, get_lexer
from pptx_utils import replace_with_image

class SlideContent:
    pass


@dataclass
class TextFormatting:
    bold: bool = False
    italic: bool = False
    underline: bool = False
    strikethrough: bool = False
    code: bool = False
    is_dark: bool = True
    
    def copy_with_changes(self, **changes):
        return TextFormatting(
            bold=changes.get('bold', self.bold),
            italic=changes.get('italic', self.italic),
            underline=changes.get('underline', self.underline),
            strikethrough=changes.get('strikethrough', self.strikethrough),
            code=changes.get('code', self.code),
            is_dark=changes.get('is_dark', self.is_dark),
        )
        
    def update(self, new_font):
        new_font.bold = self.bold
        new_font.italic = self.italic
        new_font.underline = self.underline
        new_font.strikethrough = self.strikethrough
        if self.is_dark:
            new_font.color.theme_color = MSO_THEME_COLOR.DARK_1
        else:
            new_font.color.theme_color = MSO_THEME_COLOR.LIGHT_1
        new_font.size = Pt(16)
        if self.code:
            new_font.name = "Courier New"
        return new_font
    
    def to_short_preview(self):
        return self.to_preview()
    
    def to_preview(self):
        if not any((self.bold, self.italic, self.underline, self.strikethrough, self.code, self.is_dark)):
            return ""
        return (f"{'B' if self.bold else ''}{'I' if self.italic else ''}"
                f"{'U' if self.underline else ''}{'S' if self.strikethrough else ''}"
                f"{'C' if self.code else ''}{'D' if self.is_dark else ''}")


@dataclass
class Run:
    text: str
    font: TextFormatting
    hyperlink: str = None
    
    def is_empty(self):
        return not self.text.strip()
    
    def count_lines(self):
        return textwrap.wrap(self.text, 30).count("\n") + 1
    
    def add_to_paragraph(self, paragraph):
        if self.text == "\n":
            return paragraph.add_line_break()
        r = paragraph.add_run()
        if self.font is not None:
            self.font.update(r.font)
        if self.hyperlink is not None:
            r.hyperlink.address = self.hyperlink
        r.text = self.text
        return r
    
    def to_preview(self):
        formatting = self.font.to_preview()
        if not formatting:
            return self.text
        return f"[{formatting}]{self.text}[/{formatting}]"
    
    def to_short_preview(self):
        return self.text[:10] + "..." if len(self.text) > 10 else self.text
    
@dataclass
class Paragraph:
    runs: list[Run]
    level: int
    font: TextFormatting = None
    line_spacing: float = None
    space_after: Length = None
    space_before: Length = None
    
    def is_empty(self):
        return not any(run.text.strip() for run in self.runs)
    
    def count_lines(self):
        return sum(run.count_lines() for run in self.runs) + 1
    
    def add_to_text_frame(self, text_frame, is_first_paragraph=True):
        if is_first_paragraph:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        if self.level is not None:
            p.level = self.level
        if self.line_spacing is not None:
            p.line_spacing = self.line_spacing
        if self.space_after is not None:
            p.space_after = self.space_after
        if self.space_before is not None:
            p.space_before = self.space_before
        if self.font is not None:
            p.font = self.font
        for run in self.runs:
            run.add_to_paragraph(p)
        return p
    
    def to_preview(self):
        return "    "*(self.level-1) + "|".join(run.to_preview() for run in self.runs)
    
    def to_short_preview(self):
        if not self.runs:
            return ""
        if len(self.runs) > 1:
            return self.runs[0].to_short_preview() + "\\n..."
        return self.runs[0].to_short_preview()
    
@dataclass
class TextFrame(SlideContent):
    paragraphs: list[Paragraph] = field(default_factory=list)
    auto_size: MSO_AUTO_SIZE = None
    margin_bottom: Length = None
    margin_left: Length = None
    margin_right: Length = None
    margin_top: Length = None
    vertical_anchor: MSO_VERTICAL_ANCHOR = None
    word_wrap: bool = True
    is_list: bool = False
    
    def is_empty(self):
        return not any(paragraph.runs for paragraph in self.paragraphs)
    
    def count_lines(self):
        return sum(paragraph.count_lines() for paragraph in self.paragraphs)
    
    def add_run(self, text, formatting):
        if isinstance(text, str):
            run = Run(text, formatting)
        else:
            run = text
        if not self.paragraphs:
            self.start_paragraph(0)
        self.paragraphs[-1].runs.append(run)
        return run
    
    def start_paragraph(self, level):
        self.paragraphs.append(Paragraph([], level))
    
    def add_to_slide(self, slide, left, top, width, height):
        text_box = slide.shapes.add_textbox(
            left, top, width, height
        )
        self.modify_text_frame(text_box.text_frame)
        return text_box
    
    def add_to_powerpoint(self, placeholder, slide, presentation, indent=False):
        text_frame = placeholder.text_frame
        if self.auto_size is not None:
            text_frame.auto_size = self.auto_size
        if self.vertical_anchor is not None:
            text_frame.vertical_anchor = self.vertical_anchor
        if self.word_wrap is not None:
            text_frame.word_wrap = self.word_wrap
        if self.margin_bottom is not None:
            text_frame.margin_bottom = self.margin_bottom
        if self.margin_left is not None:
            text_frame.margin_left = self.margin_left
        if self.margin_right is not None:
            text_frame.margin_right = self.margin_right
        if self.margin_top is not None:
            text_frame.margin_top = self.margin_top
        #if self.is_list:
        #    print([len(p.runs) for p in self.paragraphs])
        for index, paragraph in enumerate(self.paragraphs):
            if not paragraph.runs:
                continue
            if indent:
                paragraph.level += 1
            paragraph.add_to_text_frame(text_frame, not indent and index==0)
        return text_frame
    
    def to_preview(self):
        symbol = "list" if self.is_list else "text"
        return (f"[{symbol}]" +
                "\n".join(paragraph.to_preview() for paragraph in self.paragraphs) +
                f"[/{symbol}]")
        
    def to_short_preview(self):
        if not self.paragraphs:
            return ""
        return self.paragraphs[0].to_short_preview()
    
    
@dataclass
class Picture(SlideContent):
    image: str
    left: int = None
    top: int = None
    width: int = None
    height: int = None
    
    def is_empty(self):
        return not self.image.strip()
    
    def to_preview(self):
        return f"[image]{self.image}"
    
    def count_lines(self):
        return 1
    
    def add_to_powerpoint(self, placeholder, slide, presentation):
        replace_with_image(self.image, placeholder, slide)
        
    def to_short_preview(self):
        return self.image[:100] + "..." if len(self.image) > 100 else self.image
    
    
@dataclass
class CodeBlock(SlideContent):
    code: str
    language: str
    runnable: bool = True
    
    MAX_REASONABLE_LINE = 20
    
    def is_empty(self):
        return not self.code.strip()
    
    def count_lines(self):
        return self.code.strip().count("\n") + 1
    
    def add_to_powerpoint(self, placeholder, slide, presentation, indent=False):
        text_frame = placeholder.text_frame
        code = self.code
        # TODO: Explore possible options
        options = {}
        lexer = get_lexer(self.language, code)

        # TODO: Make this way smarter
        # if code.count("\n") < self.MAX_REASONABLE_LINE:
        formatter = PowerPointCodeFormatter(text_frame, code, overwrite=False, **options)
        result = highlight(code, lexer, formatter)
        #else:
        #    result = format_code_as_image(presentation, slide, placeholder, code, lexer, **options)
        
    def to_preview(self):
        symbol = "code" if self.runnable else "code_no_run"
        return f"[{symbol}]{self.language!r}\n{self.code}[/{symbol}]"
    
    def to_short_preview(self):
        return self.code[:30] + "..." if len(self.code) > 30 else self.code
    

@dataclass
class RawSlide:
    title: TextFrame
    type: str
    key_idea: str = None
    content: list[Union[TextFrame]] = field(default_factory=list)
    
    def copy_without_content(self):
        return RawSlide(self.title, [])
    
    def to_preview(self):
        return (f"# {self.title.to_preview()}: {self.type}\n" +
                f"Key Idea: {self.key_idea.to_preview()}\n" + 
                "\n----\n".join(content.to_preview() for content in self.content) + 
                "\n\n")
