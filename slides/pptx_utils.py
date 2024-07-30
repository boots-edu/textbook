import json

# Powerpoint stuff
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Pt

# XML Handling stuff
from lxml import etree

from pygments.styles.sas import SasStyle


def replace_with_image(img, shape, slide, max_size=False, presentation=None):
    pic = slide.shapes.add_picture(img, shape.left, shape.top)

    # calculate max width/height for target size
    ratio = min(shape.width / float(pic.width), shape.height / float(pic.height))

    if max_size:
        start_of_content_area = slide.shapes.title.top + slide.shapes.title.height
        height_of_content_area = presentation.slide_height - start_of_content_area
        height_of_content_area -= Inches(0.5)
        pic.width = int(height_of_content_area * pic.width / pic.height)
        pic.height = int(height_of_content_area)
        pic.top = start_of_content_area
        pic.left = shape.left
    else:
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        pic.top = shape.top + ((shape.height - pic.height) // 2)
        pic.left = shape.left + ((shape.width - pic.width) // 2)

    placeholder = shape.element
    placeholder.getparent().remove(placeholder)
    return


def no_bullet(paragraph):
    paragraph._pPr.insert(
        0,
        etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"),
    )


def _parse_extras(line):
    if not line:
        return {}
    return {k: json.loads(v) for part in line.split(",") for k, v in [part.split("=")]}


def format_run(ttype, run):
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    color, bold, italic, underline, _, border, roman, sans, mono = SasStyle._styles[
        ttype
    ]

    if color:
        if len(color) == 3:
            color = color[0] * 2, color[1] * 2, color[2] * 2
        else:
            color = color[:2], color[2:4], color[4:]
        color = [int(c, 16) for c in color]
        run.font.color.rgb = RGBColor(*color)
