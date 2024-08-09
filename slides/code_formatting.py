import io

# Powerpoint stuff
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Pt

# Code highlighting stuff
from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import html
from pygments.formatter import Formatter
from pygments.formatters.img import ImageFormatter
from pygments.formatters.html import HtmlFormatter
from pygments.util import ClassNotFound
from pygments.style import Style
from pygments import token

from pptx_utils import no_bullet, format_run, replace_with_image


class CodeStyle(Style):
    default_style = ""
    """
    light_mode_styles = {
        token.Whitespace: "#bbbbbb",
        token.Comment: "#008800",
        token.String: "#800080",
        token.Number: "#2c8553",
        token.Other: "bg:#ffffe0",
        token.Keyword: "#2c2cff",
        token.Keyword.Reserved: "#353580",
        token.Keyword.Constant: "",
        token.Name.Builtin: "#2c2cff",
        token.Name.Variable: "#2c2cff",
        token.Generic: "#2c2cff",
        token.Generic.Emph: "#008800",
        token.Generic.Error: "#d30202",
        token.Error: "bg:#e3d2d2 #a61717",
    }
    """
    # Dark mode
    styles = {
        token.Whitespace: "#000000",
        token.Comment: "#008800",
        token.String: "#800080",
        token.Number: "#2c8553",
        token.Other: "bg:#ffffe0",
        token.Keyword: "#2c2cff",
        token.Keyword.Reserved: "#353580",
        token.Keyword.Constant: "",
        token.Name.Builtin: "#2c2cff",
        token.Name.Variable: "#2c2cff",
        token.Generic: "#b0b0b0",
        token.Generic.Emph: "#c0c0c0",
        token.Generic.Error: "#d30202",
        token.Error: "bg:#e3d2d2 #a61717",
    }
    

class PowerPointCodeFormatter(Formatter):

    def __init__(self, text_frame, code, **options):
        self.options = options
        self.text_frame = text_frame
        self.code = code
        self.line_count = code.count("\n")
        # 9 fits comfortably

    def fix_height(self):
        paragraph = self.text_frame.paragraphs[0]
        # if self.line_count > 9:
        #    spacing = (self.MAX_REASONABLE_LINE-(self.line_count-9))/self.MAX_REASONABLE_LINE/2
        #    paragraph.line_spacing = spacing
        paragraph.line_spacing = 0.5

    def format(self, tokensource, outfile):
        if self.text_frame:
            self.text_frame.clear()
            self.text_frame.word_wrap = False
            self.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = self.text_frame.paragraphs[0]
            no_bullet(paragraph)
            paragraph.font.name = "Courier New"
            paragraph.font.size = Pt(10)
            for ttype, value in tokensource:
                run = paragraph.add_run()
                run.text = value
                format_run(ttype, run)
            self.fix_height()
        else:
            print("Throwing away codeblock!")


def format_code_as_image(presentation, current_slide, placeholder, code, lexer, **options):
    formatter = ImageFormatter(
                line_numbers=False,
                style=CodeStyle,  # get_style_by_name('sas'),
                font_size=12,
                image_pad=0,
                line_pad=8,
                **options,
            )
    image_information = highlight(code, lexer, formatter)
    with io.BytesIO() as temporary_image:
        temporary_image.write(image_information)
        temporary_image.seek(0)
        replace_with_image(
            temporary_image,
            placeholder,
            current_slide,
            True,
            presentation,
        )

def get_lexer(lang, code):
    if lang:
        try:
            return get_lexer_by_name(lang, stripall=True)
        except ClassNotFound:
            return guess_lexer(code)
    else:
        return guess_lexer(code)
    
