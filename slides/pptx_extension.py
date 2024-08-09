from pprint import pprint
import json
import math
import time
import os
from pathlib import Path
import io
from urllib.parse import unquote
from dataclasses import dataclass

# Markdown parsing stuff
import marko
from marko.ext.gfm.renderer import GFMRendererMixin
from marko.ext.gfm import elements
import marko.renderer

# Code highlighting stuff
from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import html
from pygments.formatter import Formatter
from pygments.formatters.img import ImageFormatter
from pygments.formatters.html import HtmlFormatter
from pygments.util import ClassNotFound
from pygments.style import Style
from pygments import token

# from pygments.token import Token, Comment, Keyword, Name, String, \
#    Error, Generic, Number, Operator, Whitespace
from pygments.token import (
    Keyword,
    Name,
    Comment,
    String,
    Error,
    Number,
    Operator,
    Generic,
    Literal,
)
from pygments.styles.sas import SasStyle
from pygments.styles import get_style_by_name

# Powerpoint stuff
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Pt

# XML Handling stuff
from lxml import etree

# Monkey Patches
import python_pptx_patches

# Other internal tools
from pptx_utils import replace_with_image
from code_formatting import PowerPointCodeFormatter, format_code_as_image, get_lexer
from pptx_data import Picture, RawSlide, SlideContent, Run, Paragraph, TextFrame, CodeBlock, TextFormatting


MODES = {
    "start": 0,
    "title": 1,
    "preamble": 2,
    "key idea": 3,
    "content": 4,
    "next step": 5
}

LEFT_SIDE = 1
RIGHT_SIDE = 2

GRAPHICS_FOLDER = "../source/assets/"


class SlideConverter:
    BASE_PRESENTATION = "./template-silber.pptx"
    SLIDE_LAYOUT_TYPES = {
        "title": 0,
        "title_content": 1,
        "section": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "captioned_content": 7,
        "captioned_picture": 8,
    }
    
    def __init__(self, options):
        self.options = options
        self.presentation = Presentation(self.BASE_PRESENTATION)
        
        self.current_slide = None
        self.current_text = None
        self.current_paragraph = None
        self._current_title = None
        self._current_side = LEFT_SIDE
        
    def convert(self, raw_slides):
        for slide in raw_slides:
            # Chunk up content into pairs
            content = list(slide.content)
            while content:
                # Make an appropriate slide with the right type
                if slide.type == "title":
                    current_slide = self.add_slide(slide.type)
                elif len(content) == 1:
                    current_slide = self.add_slide("captioned_content")
                else:
                    current_slide = self.add_slide("two_content")
                
                # Fill in the title
                self.add_content(current_slide, current_slide.shapes.title, slide.title)
                
                #for shape in current_slide.placeholders:
                #    print('%d %s' % (shape.placeholder_format.idx, shape.name))
                
                # Title slides only have their body
                if slide.type == "title":
                    self.add_content(current_slide, current_slide.shapes.placeholders[1], content.pop(0))
                    continue
                
                # Fill in the key idea, if available
                if slide.key_idea:
                    self.add_content(current_slide, current_slide.shapes.placeholders[13], slide.key_idea)
                
                # Fill up the content, two at a time
                if len(content) == 1:
                    self.add_content(current_slide, current_slide.shapes.placeholders[1], content.pop(0))
                else:
                    self.add_content(current_slide, current_slide.shapes.placeholders[1], content.pop(0))
                    self.add_content(current_slide, current_slide.shapes.placeholders[2], content.pop(0))
                    
    def add_content(self, slide, where, content):
        content.add_to_powerpoint(where, slide, self.presentation)
    
    def add_slide(self, type="title"):
        if type not in self.SLIDE_LAYOUT_TYPES:
            raise Exception(f"Invalid slide type: {type}")
        slide_layout_index = self.SLIDE_LAYOUT_TYPES[type]
        slide_layout = self.presentation.slide_layouts[slide_layout_index]
        current_slide = self.presentation.slides.add_slide(slide_layout)
        # current_slide.shapes[LEFT_SIDE].text_frame
        return current_slide
        

class FormattingContext:            
    def __init__(self):
        self.stack = [TextFormatting()]
        
    def __call__(self, **kwargs):
        new_context = self.current.copy_with_changes(**kwargs)
        self.stack.append(new_context)
        return self
        
    def __enter__(self):
        return self
        
    def __exit__(self, *args):
        self.stack.pop()
        return False
    
    @property
    def current(self):
        return self.stack[-1]

class PowerPointRenderer(GFMRendererMixin):
    options = {}
    # TODO: Fix these to be instance locals instead of class locals!
    

    def __init__(self, **options):
        super().__init__(**options)
        #: Estimated slide content so far, in short form
        self.mode = "start"
        self.slides: list[RawSlide] = []
        self._list = []
        self._header = None
        self._key_idea = None
        self.formatting = FormattingContext()
        
    @property
    def current_slide(self):
        return self.slides[-1] if self.slides else None
    
    @property
    def current_content(self):
        if not self.current_slide:
            return None
        if not self.current_slide.content:
            return None
        return self.current_slide.content[-1]
        
    def add_run(self, text: str) -> Run:
        if self._header:
            target = self._header
        elif self._key_idea:
            target = self._key_idea
        else:
            target = self.current_content
            if not isinstance(self.current_content, TextFrame):
                self.add_content(TextFrame())
        run = Run(text, self.formatting.current)
        target.add_run(run, self.formatting.current)
        return run
    
    def add_content(self, content: SlideContent):
        if self.current_slide is None:
            raise Exception("No slide to add content to:", content)
        self.current_slide.content.append(content)
        return content
    
    def add_slide(self, title, type="title"):
        key_idea = None
        if self.current_slide is not None:
            key_idea = self.current_slide.key_idea
        title_text_frame = TextFrame()
        title_text_frame = title_text_frame.add_run(title, self.formatting.current)
        new_slide = RawSlide(title_text_frame, type, key_idea=key_idea)
        self.slides.append(new_slide)
        return new_slide

    def start_paragraph(self, level):
        if not isinstance(self.current_content, TextFrame) or self.current_content.is_list:
            self.add_content(TextFrame())
        self.current_content.start_paragraph(level)
        
    def indent_list(self, level):
        self.current_content.start_paragraph(level)
        
    ### Rendering

    def render_strong_emphasis(self, element: "inline.StrongEmphasis") -> str:
        with self.formatting(bold=True):
            text = self.render_children(element)
        return text

    def render_emphasis(self, element: "inline.Emphasis") -> str:
        with self.formatting(italic=True):
            text = self.render_children(element)
        return text

    def render_code_span(self, element: "inline.CodeSpan") -> str:
        with self.formatting(code=True):
            text = self.render_raw_text(element)
            #text = element.children
            #self.add_run(text)
        return text
            
    #def render_text(self, element: "inline.Text") -> str:
    #    return self.add_run(element.children, self.formatting.current)
    
    def render_raw_text(self, element: "inline.RawText") -> str:
        text = element.children
        if text.startswith("{:") and text.endswith("}"):
            # TODO: Directives for no-run, boxes, etc.
            if ".no-run" in text:
                self.current_slide.content[-2].runnable = False
            return ""
        self.add_run(text)
        return text

    def render_fenced_code(self, element):
        code = element.children[0].children
        code_block = CodeBlock(code, element.lang)
        self.add_content(code_block)
        return code

    def render_heading(self, element: "block.Heading") -> str:
        header = self._header = TextFrame()
        with self.formatting(is_dark=False):
            title_text = self.render_children(element)
        self._header = None
        
        # FSA to determine what we are currently processing
        if self.mode == "start":
            self.mode = "title"
        elif self.mode == "title":
            if title_text.strip().lower() in ("key idea", "key ideas"):
                self.mode = "key idea"
            else:
                self.mode = "content"
        elif self.mode == "key idea":
            self.mode = "content"
        elif self.mode == "content":
            if title_text.strip().lower() in ("next step", "next steps"):
                self.mode = "next step"
        
        # Stop if we're not supposed to grab this title        
        if self.mode in ("start", "preamble", "next step", "key idea"):
            return ""
        
        # Create a new slide, either a title or a two_content
        if element.level == 1:
            new_slide = self.add_slide(header, "title")
            # TODO: Extract this to a setting!
            content = TextFrame()
            self.add_content(content)
            content.add_run("BOOTS: Beginner Object-Oriented TypeScript", self.formatting.current)
        else:
            new_slide = self.add_slide(header, "two_content")
        
        return title_text
    
    IGNORED_CONTENT_MODES = ("start", "title", "preamble", "next step")

    def render_paragraph(self, element: "block.Paragraph") -> str:
        # Abort if we're not supposed to be rendering this
        if self.mode in self.IGNORED_CONTENT_MODES:
            return ""
        
        # TODO: Stop if we're rendering a class attribute
        #if text.startswith("{:") and text.endswith("}"):
        #    return ""
        if self._list:
            text = self.render_children(element)
        # Key idea gets saved
        elif self.mode == "key idea":
            # print("SEEN KEY IDEA", element)
            key_idea = self._key_idea = TextFrame()
            with self.formatting(is_dark=False):
                text = self.render_children(element)
            self._key_idea = None
            self.current_slide.key_idea = key_idea
        else:
            self.start_paragraph(0)
            text = self.render_children(element)
            
        return text
    
    def render_line_break(self, element: "inline.LineBreak") -> str:
        self.add_run(" ")
        return " "
        
    def render_list(self, element: "block.List") -> str:
        if self.mode in self.IGNORED_CONTENT_MODES:
            return ""
        
        if self._list:
            self.indent_list(len(self._list)-1)
        else:
            self.add_content(TextFrame(is_list=True))
        self._list.append(element)
        children = self.render_children(element)
        self._list.pop()
        
        return children

    def render_list_item(self, element: "block.ListItem") -> str:
        self.indent_list(len(self._list)-1)
        text = self.render_children(element)
        return text

    def render_image(self, element: "inline.Image") -> str:
        if self.mode in self.IGNORED_CONTENT_MODES:
            return ""
        
        image_path = Path(element.dest).parts[2:]
        image_path = os.path.join(*image_path)
        url = unquote(self.escape_url(os.path.join(self.GRAPHICS_FOLDER, image_path)))
        
        self.add_content(Picture(url))
        
        return url

    def finish(self):
        # Pretty print self.slides (list of dataclasses)
        # content = "#".join(slide.to_preview() for slide in self.slides)
        # with open("dist/var_slides.txt", "w") as f:
        #     print(content, file=f)
        return ""

    @staticmethod
    def escape_html(raw: str) -> str:
        return raw

    @staticmethod
    def escape_url(raw: str) -> str:
        return raw
    
    def render_table(self, element: "block.Table") -> str:
        raise NotImplementedError("Tables are not yet supported in PowerPoint output")
    
    def render_url(self, element: "inline.Url") -> str:
        raise NotImplementedError("URLs are not yet supported in PowerPoint output")
    
    def render_quote(self, element: "block.Quote") -> str:
        # TODO: Finish this
        return self.render_children(element)
        #return self.render_paragraph(element)
    
    def render_thematic_break(self, element: "block.ThematicBreak") -> str:
        raise NotImplementedError("Thematic breaks are not yet supported in PowerPoint output")


class PPTXRenderExtension:
    elements = [
        elements.Paragraph,
        elements.Strikethrough,
        elements.Url,
        elements.Table,
        elements.TableRow,
        elements.TableCell,
    ]
    renderer_mixins = [PowerPointRenderer]
    parser_mixins = []


# XML Stuff

ETREE_NAMESPACE_MAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "p159": "http://schemas.microsoft.com/office/powerpoint/2015/09/main",
}

for k, v in ETREE_NAMESPACE_MAP.items():
    etree.register_namespace(k, v)


def xpath(el, query):
    return etree.ElementBase.xpath(el, query, namespaces=ETREE_NAMESPACE_MAP)
